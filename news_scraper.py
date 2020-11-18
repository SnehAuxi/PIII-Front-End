import requests
import os
from pprint import pprint
import urllib
import urllib.request as urllib2
import os
from pptx import Presentation

def news_scraper(topic):
    apikey = "OuCSULEj80JbINMUCsuKHpoMOcuilkIi"
    title = []
    body = []
    img_list = []
    query = topic
    begin_date = "20190701"  # YYYYMMDD
    page = "0"  # <0-100>
    sort = "relevance"  # newest, oldest
    query_url = f"https://api.nytimes.com/svc/search/v2/articlesearch.json?" \
                f"q={query}" \
                f"&api-key={apikey}" \
                f"&begin_date={begin_date}" \
                f"&page={page}" \
                f"&sort={sort}"

    r = requests.get(query_url)
    response = r.json()
    for i in range(0, len(response['response']['docs'])):
        title.append(response['response']['docs'][i]['headline']['main'])
        body.append(response['response']['docs'][i]['lead_paragraph'])
        if response['response']['docs'][i]['multimedia']:
            url = str('http://static01.nyt.com/' + response['response']['docs'][i]['multimedia'][0]['url'])
            urllib.request.urlretrieve(url, os.path.basename("../images/" + url))
            img_list.append("../images/" + url.split("/")[-1])
        
    return (title, body, img_list)


def alexa_pres(pres,  num_slides, headers_list, text_list: list, img_list: list,data_len=0):
    if not headers_list:
        for x in range(0,num_slides):
            slide = pres.slides.add_slide(pres.slide_layouts[8])

    elif data_len != 0:
        for x in range(0, data_len - 1):
            slide = pres.slides.add_slide(pres.slide_layouts[8])
            shapes = slide.shapes
            title = shapes.title
            title.text = headers_list[x]
            try:
                if img_list[x]:
                    picture_placeholder = slide.placeholders[1]
                    picture_placeholder = picture_placeholder.insert_picture(img_list[x])
            except:
                    print("")
            text_place = slide.placeholders[2]
            text_place.text  = text_list[x] 
        diff = num_slides-data_len
        for x in range(0, diff):
            slide = pres.slides.add_slide(pres.slide_layouts[8])
    else:
            for x in range(0, num_slides - 1):
                slide = pres.slides.add_slide(pres.slide_layouts[8])
                shapes = slide.shapes
                title = shapes.title
                title.text = headers_list[x]
                try:
                    if img_list[x]:
                        picture_placeholder = slide.placeholders[1]
                        picture_placeholder = picture_placeholder.insert_picture(img_list[x])
                except:
                        print("")
                text_place = slide.placeholders[2]
                text_place.text  = text_list[x] 
            
    pres.save('Final - Presentation ny.pptx')

def create_ppt_ny(ppt_subject, count):
    if not os.path.exists("images/"):
        os.mkdir("images/")
    os.chdir("images/")
    (title, body, img_list) = news_scraper(ppt_subject)
    os.chdir(os.path.abspath(os.path.join(os.getcwd(), os.pardir)))
    
    pres =  Presentation()
    if len(title)<count and len(title) != 0:
        alexa_pres(pres, count, title, body, img_list,len(title))    
    else:
        alexa_pres(pres, count, title, body, img_list)
    return 'Created the presentation from New York Times'
